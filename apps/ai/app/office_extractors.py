"""Extractores de documentos de oficina: DOCX, XLSX, XLS y PPTX.

Portados del módulo de lectura de archivos de Paloma
(paloma-os/core/file_processing.py) y adaptados al contrato de
``files.extract_text``: cada extractor devuelve SIEMPRE un ``str`` (el texto
extraído, o una nota entre corchetes si no se pudo), y NUNCA lanza. Las
dependencias (python-docx, openpyxl, xlrd, python-pptx) se importan de forma
perezosa para que importar este módulo no falle si alguna no está instalada.
"""
from __future__ import annotations

import logging
from pathlib import Path

log = logging.getLogger("seguria.office")

MAX_TEXT = 15_000  # chars que se entregan al LLM por documento
_ROWS_PER_SHEET = 200  # cap de filas serializadas por hoja de Excel


def _cap(text: str) -> str:
    text = (text or "").strip()
    if len(text) > MAX_TEXT:
        return text[:MAX_TEXT] + "\n\n[... texto truncado ...]"
    return text


def extract_docx(path: str) -> str:
    """Texto de un DOCX (párrafos no vacíos) usando python-docx."""
    try:
        from docx import Document  # import perezoso

        doc = Document(str(path))
        parts = [p.text for p in doc.paragraphs if p.text.strip()]
        # Tablas: cada fila como `celda | celda | ...`
        for table in doc.tables:
            for row in table.rows:
                cells = [c.text.strip() for c in row.cells]
                if any(cells):
                    parts.append(" | ".join(cells))
        return _cap("\n".join(parts)) or f"[Documento Word sin texto extraíble: {Path(path).name}]"
    except ImportError:
        return f"[Documento Word adjunto: {Path(path).name} — instalar python-docx para leerlo]"
    except Exception as exc:
        log.warning("DOCX %s: %s", path, exc)
        return f"[Documento Word adjunto: {Path(path).name} — no se pudo extraer el texto]"


def extract_xlsx(path: str) -> str:
    """Texto de un XLSX: filas no vacías por hoja como `valor | valor | ...`."""
    try:
        from openpyxl import load_workbook  # import perezoso

        wb = load_workbook(path, read_only=True, data_only=True)
        parts: list[str] = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            parts.append(f"## Hoja: {sheet_name}")
            lines: list[str] = []
            for row in ws.iter_rows(values_only=True):
                cells = [str(c) if c is not None else "" for c in row]
                if not any(cells):
                    continue
                lines.append(" | ".join(cells))
                if len(lines) >= _ROWS_PER_SHEET:
                    lines.append("[... filas truncadas ...]")
                    break
            if lines:
                parts.append("\n".join(lines))
        wb.close()
        return _cap("\n\n".join(parts)) or f"[Excel sin contenido legible: {Path(path).name}]"
    except ImportError:
        return f"[Excel adjunto: {Path(path).name} — instalar openpyxl para leerlo]"
    except Exception as exc:
        log.warning("XLSX %s: %s", path, exc)
        return f"[Excel adjunto: {Path(path).name} — no se pudo extraer el contenido]"


def extract_xls(path: str) -> str:
    """Excel legacy (.xls binario) con xlrd; degrada a nota si no está."""
    try:
        import xlrd  # type: ignore  # import perezoso

        wb = xlrd.open_workbook(str(path))
        parts: list[str] = []
        for idx in range(wb.nsheets):
            sh = wb.sheet_by_index(idx)
            parts.append(f"## Hoja: {sh.name}")
            lines = []
            for r in range(min(sh.nrows, _ROWS_PER_SHEET)):
                cells = [str(sh.cell_value(r, c)) for c in range(sh.ncols)]
                if any(cells):
                    lines.append(" | ".join(cells))
            if lines:
                parts.append("\n".join(lines))
        return _cap("\n\n".join(parts)) or f"[Excel legacy sin contenido: {Path(path).name}]"
    except ImportError:
        return (f"[Excel legacy adjunto: {Path(path).name} — formato .xls requiere xlrd; "
                "pídele al cliente convertirlo a .xlsx]")
    except Exception as exc:
        log.warning("XLS %s: %s", path, exc)
        return f"[Excel legacy adjunto: {Path(path).name} — no se pudo extraer]"


def extract_pptx(path: str) -> str:
    """Texto de un PPTX por slide (incluye notas del orador)."""
    try:
        from pptx import Presentation  # import perezoso

        prs = Presentation(str(path))
        slides: list[str] = []
        for i, slide in enumerate(prs.slides, start=1):
            chunks: list[str] = []
            for shape in slide.shapes:
                if getattr(shape, "has_text_frame", False):
                    for p in shape.text_frame.paragraphs:
                        line = "".join(run.text for run in p.runs).strip()
                        if line:
                            chunks.append(line)
            try:
                if slide.has_notes_slide:
                    notes = slide.notes_slide.notes_text_frame.text.strip()
                    if notes:
                        chunks.append(f"[notas] {notes}")
            except Exception:
                pass
            if chunks:
                slides.append(f"## Slide {i}\n" + "\n".join(chunks))
        return _cap("\n\n".join(slides)) or f"[Presentación sin texto extraíble: {Path(path).name}]"
    except ImportError:
        return f"[Presentación adjunta: {Path(path).name} — instalar python-pptx para leerla]"
    except Exception as exc:
        log.warning("PPTX %s: %s", path, exc)
        return f"[Presentación adjunta: {Path(path).name} — no se pudo extraer el texto]"
